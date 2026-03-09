#!/usr/bin/env python3
"""
Dropbox Indexer v5 — Delta-based: only fetch CHANGED files since last run.
Uses Dropbox list_folder/continue with saved cursor for incremental updates.
"""

import json, os, re, time, io, sys, urllib.request, urllib.error, subprocess, tempfile
from pathlib import Path
from datetime import datetime

# ── Config ─────────────────────────────────────────────────────────────────
ENV_FILE        = os.path.expanduser("~/.openclaw/.env")
OPENCLAW_HOME   = Path(os.environ.get("OPENCLAW_HOME", os.path.expanduser("~/.openclaw")))
OUTPUT_DIR      = OPENCLAW_HOME / "workspace/memory/knowledge/dropbox"
PROGRESS_FILE   = OPENCLAW_HOME / "workspace/memory/dropbox-index-progress.json"
CURSOR_FILE     = OPENCLAW_HOME / "workspace/memory/dropbox-cursor.json"
LOG_FILE        = OPENCLAW_HOME / "workspace/memory/dropbox-indexer.log"

FOLDERS = ["/Documents", "/Work", "/Research"]
SKIP_PATHS = ["/Archive", "/Backups", "/Photos"]

INDEX_EXTS = {"pdf", "docx", "doc", "txt", "md", "markdown", "rst", "csv", "json", "tex", "jpg", "jpeg", "png", "xlsx", "xls", "pptx", "ppt"}

MAX_FILE_SIZE = 20 * 1024 * 1024
MAX_TEXT_CHARS = 12000

# ── Helpers ─────────────────────────────────────────────────────────────────
def load_env():
    env = {}
    with open(ENV_FILE) as f:
        for line in f:
            line = line.strip()
            if "=" in line and not line.startswith("#"):
                k, v = line.split("=", 1)
                env[k.strip()] = v.strip()
    return env

def get_token(env):
    data = (
        f"grant_type=refresh_token"
        f"&refresh_token={env['DROPBOX_FULL_REFRESH_TOKEN']}"
        f"&client_id={env['DROPBOX_FULL_APP_KEY']}"
        f"&client_secret={env['DROPBOX_FULL_APP_SECRET']}"
    ).encode()
    req = urllib.request.Request(
        "https://api.dropboxapi.com/oauth2/token", data=data,
        headers={"Content-Type": "application/x-www-form-urlencoded"}
    )
    return json.loads(urllib.request.urlopen(req).read())["access_token"]

def api_call(token, endpoint, data, retries=5):
    url = f"https://api.dropboxapi.com/2/files/{endpoint}"
    for attempt in range(retries):
        try:
            req = urllib.request.Request(
                url, data=json.dumps(data).encode(),
                headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
            )
            return json.loads(urllib.request.urlopen(req, timeout=30).read())
        except urllib.error.HTTPError as e:
            if e.code == 429:
                wait = 2 ** attempt * 5
                log(f"  ⏳ Rate limited, waiting {wait}s...")
                time.sleep(wait)
            elif e.code in (500, 503):
                time.sleep(2 ** attempt)
            else:
                return {}
        except Exception:
            time.sleep(2 ** attempt)
    return {}

def download_file(token, path, retries=4):
    for attempt in range(retries):
        try:
            req = urllib.request.Request(
                "https://content.dropboxapi.com/2/files/download",
                headers={"Authorization": f"Bearer {token}",
                         "Dropbox-API-Arg": json.dumps({"path": path})}
            )
            return urllib.request.urlopen(req, timeout=60).read()
        except Exception:
            time.sleep(2 ** attempt)
    return None

def list_folder_delta(token, folder, cursor=None):
    """List folder changes since cursor (or all if cursor=None)."""
    entries = []
    
    if cursor:
        # Continue from cursor (only changes)
        r = api_call(token, "list_folder/continue", {"cursor": cursor})
    else:
        # Initial listing
        r = api_call(token, "list_folder", {"path": folder, "recursive": True, "limit": 2000})
    
    entries.extend(r.get("entries", []))
    new_cursor = r.get("cursor")
    
    while r.get("has_more"):
        time.sleep(0.5)
        r = api_call(token, "list_folder/continue", {"cursor": r["cursor"]})
        entries.extend(r.get("entries", []))
        new_cursor = r.get("cursor")
    
    return entries, new_cursor

# Text extraction functions (same as v4)
def extract_text_pdf(data):
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        text = "".join(page.extract_text() or "" for page in reader.pages[:30])
        if len(text.strip()) > 100:
            return text.strip()
    except Exception:
        pass
    return ocr_pdf(data)

def ocr_pdf(data):
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(data)
            tmp = f.name
        parts = []
        subprocess.run(["pdftoppm", "-r", "150", "-l", "5", tmp, "/tmp/ocr_p"],
                       capture_output=True, timeout=60)
        for img in sorted(Path("/tmp").glob("ocr_p*")):
            r = subprocess.run(["tesseract", str(img), "stdout", "-l", "eng+deu"],
                               capture_output=True, text=True, timeout=30)
            parts.append(r.stdout)
            img.unlink(missing_ok=True)
        os.unlink(tmp)
        return "\n".join(parts).strip()
    except Exception as e:
        return f"[OCR failed: {e}]"

def ocr_image(data):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as f:
            f.write(data)
            tmp = f.name
        r = subprocess.run(["tesseract", tmp, "stdout", "-l", "eng+deu"],
                           capture_output=True, text=True, timeout=30)
        os.unlink(tmp)
        return r.stdout.strip()
    except Exception as e:
        return f"[OCR failed: {e}]"

def extract_text_docx(data):
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[Error: {e}]"

def extract_text_xlsx(data):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets[:5]:
            parts.append(f"## Sheet: {sheet.title}")
            for row in list(sheet.iter_rows(values_only=True))[:100]:
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as e:
        return f"[Excel extraction failed: {e}]"

def extract_text_pptx(data):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides[:30], 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[PowerPoint extraction failed: {e}]"

def extract_text(path, data):
    ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""
    try:
        if ext == "pdf":
            return extract_text_pdf(data)
        elif ext in ("docx", "doc"):
            return extract_text_docx(data)
        elif ext in ("xlsx", "xls"):
            return extract_text_xlsx(data)
        elif ext in ("pptx", "ppt"):
            return extract_text_pptx(data)
        elif ext in ("jpg", "jpeg", "png"):
            return ocr_image(data)
        elif ext in ("txt", "md", "markdown", "rst", "csv", "json", "tex"):
            return data.decode("utf-8", errors="replace")
    except Exception as e:
        return f"[Error: {e}]"
    return None

def safe_filename(path):
    return re.sub(r"[^\w\-.]", "_", path.lower().strip("/"))[:200]

def load_progress():
    if PROGRESS_FILE.exists():
        return json.loads(PROGRESS_FILE.read_text())
    return {"indexed": {}, "skipped": [], "failed": [], "last_updated": None}

def save_progress(p):
    p["last_updated"] = datetime.utcnow().isoformat() + "Z"
    PROGRESS_FILE.write_text(json.dumps(p, indent=2))

def load_cursors():
    if CURSOR_FILE.exists():
        return json.loads(CURSOR_FILE.read_text())
    return {}

def save_cursors(cursors):
    CURSOR_FILE.write_text(json.dumps(cursors, indent=2))

def log(msg):
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    line = f"[{timestamp}] {msg}\n"
    print(msg)
    with open(LOG_FILE, "a") as f:
        f.write(line)

# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    env = load_env()
    log("🔑 Getting Dropbox token...")
    token = get_token(env)
    
    progress = load_progress()
    indexed_hashes = progress["indexed"]  # {path: content_hash}
    cursors = load_cursors()

    log(f"📋 Delta scan (already indexed: {len(indexed_hashes)})...")
    
    all_changes = []
    for folder in FOLDERS:
        cursor = cursors.get(folder)
        if cursor:
            log(f"  Checking {folder} for changes (delta mode)...")
        else:
            log(f"  Initial listing: {folder}...")
        
        entries, new_cursor = list_folder_delta(token, folder, cursor)
        all_changes.extend(entries)
        cursors[folder] = new_cursor
        log(f"    {len(entries)} entries")
    
    save_cursors(cursors)
    
    # Filter to new/changed indexable files
    to_index = []
    skipped = 0
    deleted = 0
    updated = 0

    for e in all_changes:
        path = e.get("path_display") or e.get("path_lower", "")

        # Handle deletions
        if e.get(".tag") == "deleted":
            if path in indexed_hashes:
                del indexed_hashes[path]
                deleted += 1
            continue

        if e[".tag"] != "file":
            continue

        path_lower = e["path_lower"]
        ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""
        size = e.get("size", 0)
        content_hash = e.get("content_hash", "")

        # Check if already indexed with same content
        if path in indexed_hashes:
            if indexed_hashes[path] == content_hash and content_hash:
                # Content unchanged — skip
                continue
            # Content changed — re-index
            log(f"  🔄 Changed: {Path(path).name}")
            updated += 1

        # Skip by path
        if any(skip in path_lower for skip in SKIP_PATHS):
            skipped += 1
            continue

        # Skip by extension
        if ext not in INDEX_EXTS:
            skipped += 1
            continue

        # Skip if too large
        if size > MAX_FILE_SIZE:
            log(f"  ⏭️ Skipping (too large): {path} ({size//1024}KB)")
            skipped += 1
            continue

        to_index.append({"path": path, "size": size, "ext": ext, "content_hash": content_hash})

    log(f"✅ {len(to_index)} to index ({updated} updated), {skipped} skipped, {deleted} deleted")
    
    if not to_index and deleted == 0:
        log("✨ Nothing changed. All done!")
        save_progress(progress)
        return
    
    # Index new files
    new_indexed = 0
    new_failed = 0
    
    for i, file_info in enumerate(to_index, 1):
        path = file_info["path"]
        content_hash = file_info.get("content_hash", "")
        log(f"[{i}/{len(to_index)}] {Path(path).name}...")
        
        data = download_file(token, path)
        if not data:
            log(f"  ❌ Download failed")
            progress["failed"].append(path)
            new_failed += 1
            continue
        
        text = extract_text(path, data)
        if not text or len(text.strip()) < 20:
            log(f"  ⏭️ No text extracted")
            progress["skipped"].append(path)
            continue
        
        text = text[:MAX_TEXT_CHARS]
        out_name = safe_filename(path) + ".md"
        header = f"# {Path(path).name}\n**Path:** {path}\n**Extracted:** {datetime.now().strftime('%Y-%m-%d')}\n\n"
        (OUTPUT_DIR / out_name).write_text(header + text)
        
        indexed_hashes[path] = content_hash
        new_indexed += 1
        
        if i % 10 == 0:
            save_progress(progress)
        
        time.sleep(0.3)
    
    save_progress(progress)
    log(f"🎉 Done: {new_indexed} indexed, {new_failed} failed, {deleted} deleted")

if __name__ == "__main__":
    main()
